import pandas as pd
import pdfplumber
from PIL import Image
import pytesseract
from pptx import Presentation
import io
from typing import Dict, List, Any
import numpy as np
from langchain_core.documents import Document
import os
import streamlit as st
import tempfile
from langchain.document_loaders import PyPDFLoader, UnstructuredPDFLoader
import zipfile
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
import shutil
import webvtt
import pysrt
from moviepy.editor import VideoFileClip
import speech_recognition as sr
from pydub import AudioSegment

def handle_pdf(file) -> tuple[List[Document], List[Dict]]:
    """Handle large PDF files with multiple extraction methods"""
    documents = []
    tables = []
    
    try:
        # First try with pdfplumber
        with pdfplumber.open(file) as pdf:
            total_pages = len(pdf.pages)
            with st.spinner(f"Processing {total_pages} pages..."):
                for page_num in range(total_pages):
                    try:
                        page = pdf.pages[page_num]
                        
                        # Extract text
                        text = page.extract_text()
                        if text and text.strip():
                            cleaned_text = text.strip().replace('\n\n', '\n').replace('\x00', '')
                            if cleaned_text:  # Verify cleaned text is not empty
                                documents.append(Document(
                                    page_content=cleaned_text,
                                    metadata={
                                        "source": file.name,
                                        "page": page_num + 1,
                                        "type": "pdf_text"
                                    }
                                ))
                        
                        # Extract tables
                        page_tables = page.extract_tables()
                        for table_num, table in enumerate(page_tables, 1):
                            if table and any(any(cell for cell in row) for row in table):
                                # Clean and validate table data
                                cleaned_table = [
                                    [str(cell).strip() if cell else "" for cell in row]
                                    for row in table if any(cell for cell in row)  # Skip empty rows
                                ]
                                
                                if cleaned_table:  # Only add if table has content
                                    # Convert table to text format for better indexing
                                    table_text = f"Table {table_num} on page {page_num + 1}:\n"
                                    for row in cleaned_table:
                                        table_text += " | ".join(row) + "\n"
                                    
                                    if table_text.strip():  # Verify table text is not empty
                                        # Add table as document for text search
                                        documents.append(Document(
                                            page_content=table_text,
                                            metadata={
                                                "source": file.name,
                                                "page": page_num + 1,
                                                "type": "pdf_table",
                                                "table_id": f"table_{page_num + 1}_{table_num}"
                                            }
                                        ))
                                        
                                        # Store raw table data
                                        tables.append({
                                            "raw_data": cleaned_table,
                                            "page": page_num + 1,
                                            "source": file.name,
                                            "type": "pdf_table",
                                            "table_id": f"table_{page_num + 1}_{table_num}"
                                        })
                    except Exception as e:
                        st.warning(f"Error processing page {page_num + 1}: {str(e)}")
        
        # If no content was extracted, try PyPDFLoader as fallback
        if not documents and not tables:
            with st.spinner("Trying alternative extraction method..."):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                    tmp_file.write(file.getvalue())
                    tmp_file_path = tmp_file.name
                
                try:
                    loader = PyPDFLoader(tmp_file_path)
                    fallback_docs = loader.load()
                    if fallback_docs:
                        # Clean and process the extracted text
                        for doc in fallback_docs:
                            if doc.page_content and doc.page_content.strip():
                                cleaned_text = doc.page_content.strip().replace('\n\n', '\n').replace('\x00', '')
                                if cleaned_text:  # Verify cleaned text is not empty
                                    documents.append(Document(
                                        page_content=cleaned_text,
                                        metadata={
                                            "source": file.name,
                                            "page": doc.metadata.get("page", 0),
                                            "type": "pdf_text"
                                        }
                                    ))
                except Exception as e:
                    st.warning(f"Alternative extraction failed: {str(e)}")
                finally:
                    if os.path.exists(tmp_file_path):
                        os.remove(tmp_file_path)
    
    except Exception as e:
        st.error(f"Error processing PDF: {str(e)}")
        # Try one last time with UnstructuredPDFLoader
        with st.spinner("Attempting final extraction method..."):
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                    tmp_file.write(file.getvalue())
                    tmp_file_path = tmp_file.name
                
                loader = UnstructuredPDFLoader(tmp_file_path)
                final_docs = loader.load()
                if final_docs:
                    # Clean and process the extracted text
                    for doc in final_docs:
                        if doc.page_content and doc.page_content.strip():
                            cleaned_text = doc.page_content.strip().replace('\n\n', '\n').replace('\x00', '')
                            if cleaned_text:  # Verify cleaned text is not empty
                                documents.append(Document(
                                    page_content=cleaned_text,
                                    metadata={
                                        "source": file.name,
                                        "page": doc.metadata.get("page_number", 0),
                                        "type": "pdf_text"
                                    }
                                ))
                
                if os.path.exists(tmp_file_path):
                    os.remove(tmp_file_path)
            except Exception as final_e:
                st.error(f"All extraction methods failed: {str(final_e)}")
    
    # Return results
    if documents:
        st.success(f"Successfully processed {len(documents)} sections from {file.name}")
        return documents, tables
    else:
        st.error("No valid content could be extracted from the PDF")
        return [], []

def handle_excel(file) -> tuple[List[Document], List[Dict]]:
    """Handle large Excel files"""
    documents = []
    tables = []
    
    try:
        # Read Excel file
        excel_file = pd.ExcelFile(file)
        total_sheets = len(excel_file.sheet_names)
        
        for sheet_idx, sheet_name in enumerate(excel_file.sheet_names, 1):
            with st.spinner(f"Processing sheet {sheet_idx}/{total_sheets}: {sheet_name}"):
                try:
                    # Read the entire sheet
                    df = pd.read_excel(
                        file,
                        sheet_name=sheet_name
                    )
                    
                    if not df.empty:
                        # Convert DataFrame to string representation
                        text_content = f"Sheet: {sheet_name}\n"
                        text_content += df.to_string(index=False)
                        
                        # Add to documents
                        documents.append(Document(
                            page_content=text_content,
                            metadata={
                                "source": file.name,
                                "sheet": sheet_name,
                                "type": "excel_text",
                                "rows": len(df)
                            }
                        ))
                        
                        # Store table data
                        tables.append({
                            "raw_data": df.values.tolist(),
                            "headers": df.columns.tolist(),
                            "sheet": sheet_name,
                            "source": file.name,
                            "type": "excel_table",
                            "rows": len(df)
                        })
                        
                except Exception as e:
                    st.warning(f"Error processing sheet {sheet_name}: {str(e)}")
                    continue
                
    except Exception as e:
        st.error(f"Error reading Excel file: {str(e)}")
        return documents, tables
    
    return documents, tables

def handle_powerpoint(file) -> List[Document]:
    """Handle PowerPoint files"""
    documents = []
    
    presentation = Presentation(file)
    for slide_num, slide in enumerate(presentation.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        
        if slide_text:
            documents.append(Document(
                page_content="\n".join(slide_text),
                metadata={
                    "source": file.name,
                    "slide": slide_num,
                    "type": "powerpoint"
                }
            ))
    
    return documents

def handle_image(file) -> tuple[List[Document], List[Dict]]:
    """Handle large images"""
    documents = []
    images = []
    
    try:
        # Open image
        with Image.open(file) as image:
            # Calculate new dimensions while maintaining aspect ratio
            max_dimension = 2000  # Maximum dimension for processing
            ratio = min(max_dimension / image.width, max_dimension / image.height)
            new_size = (int(image.width * ratio), int(image.height * ratio))
            
            # Resize image if necessary
            if ratio < 1:
                image = image.resize(new_size, Image.Resampling.LANCZOS)
            
            # Convert to RGB if necessary
            if image.mode not in ('L', 'RGB'):
                image = image.convert('RGB')
            
            # Process image in memory-efficient manner
            with st.spinner("Processing image with OCR..."):
                try:
                    # Set Tesseract path for Windows
                    if os.name == 'nt':
                        pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
                    
                    # Perform OCR with optimized settings
                    text = pytesseract.image_to_string(
                        image,
                        lang='eng',
                        config='--psm 3 --oem 3 --dpi 300'
                    ).strip()
                    
                    if text:
                        documents.append(Document(
                            page_content=text,
                            metadata={
                                "source": file.name,
                                "type": "image_text",
                                "dimensions": f"{image.width}x{image.height}"
                            }
                        ))
                    
                    # Save processed image
                    buffered = io.BytesIO()
                    image.save(buffered, format='JPEG', quality=85, optimize=True)
                    
                    images.append({
                        "image": image,
                        "text": text if text else "No text extracted",
                        "source": file.name,
                        "type": "image",
                        "dimensions": f"{image.width}x{image.height}"
                    })
                    
                except Exception as e:
                    st.warning(f"OCR processing error: {str(e)}")
                    images.append({
                        "image": image,
                        "text": "OCR processing failed",
                        "source": file.name,
                        "type": "image",
                        "dimensions": f"{image.width}x{image.height}"
                    })
    
    except Exception as e:
        raise RuntimeError(f"Error processing image: {str(e)}")
    
    return documents, images

@st.cache_data(show_spinner=False)
def extract_video_text(file_path: str) -> str:
    """Extract text from video files including audio transcription"""
    text_content = []
    
    # Check for subtitle files with the same name
    base_path = os.path.splitext(file_path)[0]
    
    # Try VTT subtitles
    vtt_path = base_path + '.vtt'
    if os.path.exists(vtt_path):
        try:
            for caption in webvtt.read(vtt_path):
                text_content.append(caption.text)
        except Exception as e:
            st.warning(f"Error reading VTT subtitles: {str(e)}")
    
    # Try SRT subtitles
    srt_path = base_path + '.srt'
    if os.path.exists(srt_path):
        try:
            subs = pysrt.open(srt_path)
            for sub in subs:
                text_content.append(sub.text)
        except Exception as e:
            st.warning(f"Error reading SRT subtitles: {str(e)}")
    
    # Process video content
    try:
        with st.spinner("Processing video content..."):
            video = VideoFileClip(file_path)
            
            # Extract video metadata
            text_content.append(f"Video Duration: {video.duration} seconds")
            text_content.append(f"Resolution: {video.size[0]}x{video.size[1]}")
            
            # Process audio if present
            if video.audio:
                text_content.append("Audio Present: Yes")
                
                try:
                    # Extract audio and convert to WAV
                    temp_audio = os.path.join(os.path.dirname(file_path), "temp_audio.wav")
                    video.audio.write_audiofile(temp_audio, codec='pcm_s16le', verbose=False, logger=None)
                    
                    # Initialize speech recognition
                    recognizer = sr.Recognizer()
                    
                    # Adjust recognition settings
                    recognizer.energy_threshold = 300
                    recognizer.dynamic_energy_threshold = True
                    recognizer.pause_threshold = 0.8
                    
                    try:
                        # Load audio file
                        with sr.AudioFile(temp_audio) as source:
                            # Adjust for ambient noise
                            recognizer.adjust_for_ambient_noise(source, duration=0.5)
                            
                            # Process audio in larger chunks for better performance
                            chunk_duration = 60  # increased from 30 to 60 seconds
                            total_duration = video.duration
                            chunks_text = []
                            
                            for i in range(0, int(total_duration), chunk_duration):
                                # Read chunk of audio
                                chunk = recognizer.record(source, duration=min(chunk_duration, total_duration - i))
                                
                                try:
                                    # Try with Google Speech Recognition
                                    chunk_text = recognizer.recognize_google(chunk, language='en-US')
                                    if chunk_text:
                                        chunks_text.append(chunk_text)
                                except sr.RequestError as e:
                                    continue
                                except sr.UnknownValueError:
                                    continue
                            
                            if chunks_text:
                                text_content.append("\nTranscribed Audio:")
                                text_content.extend(chunks_text)
                            else:
                                text_content.append("\nNote: No clear speech could be transcribed")
                            
                    except Exception as e:
                        text_content.append("\nNote: Audio processing failed")
                    
                    finally:
                        # Clean up temporary audio file
                        if os.path.exists(temp_audio):
                            os.remove(temp_audio)
                            
                except Exception as e:
                    text_content.append("\nNote: Audio extraction failed")
            
            video.close()
            
    except Exception as e:
        st.warning(f"Error processing video: {str(e)}")
    
    # Return whatever content we managed to extract
    return "\n".join(text_content) if text_content else ""

def handle_scorm(file) -> List[Document]:
    """Handle SCORM package files including video content"""
    documents = []
    
    try:
        # Create a temporary directory to extract SCORM contents
        with tempfile.TemporaryDirectory() as temp_dir:
            # First, save the uploaded file
            temp_zip = os.path.join(temp_dir, "scorm_package.zip")
            with open(temp_zip, "wb") as f:
                f.write(file.getvalue())
            
            # Extract the SCORM package
            with zipfile.ZipFile(temp_zip, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
            
            # Process manifest file
            manifest_path = os.path.join(temp_dir, "imsmanifest.xml")
            if os.path.exists(manifest_path):
                try:
                    tree = ET.parse(manifest_path)
                    root = tree.getroot()
                    # Extract metadata from manifest
                    manifest_content = ET.tostring(root, encoding='unicode', method='xml')
                    # Parse manifest for resource information
                    resources = []
                    for resource in root.findall(".//{http://www.imsproject.org/xsd/imscp_rootv1p1p2}resource"):
                        resources.append({
                            "id": resource.get("identifier"),
                            "type": resource.get("type"),
                            "href": resource.get("href")
                        })
                    
                    documents.append(Document(
                        page_content=f"{manifest_content}\nResources:\n" + 
                                   "\n".join([f"- {r['id']}: {r['type']} ({r['href']})" for r in resources]),
                        metadata={
                            "source": file.name,
                            "type": "scorm_manifest",
                            "section": "metadata",
                            "resources": resources
                        }
                    ))
                except Exception as e:
                    st.warning(f"Error processing manifest: {str(e)}")
            
            # Process all files
            for root, _, files in os.walk(temp_dir):
                for filename in files:
                    file_path = os.path.join(root, filename)
                    file_ext = os.path.splitext(filename)[1].lower()
                    
                    try:
                        # Process HTML files
                        if file_ext in ['.html', '.htm']:
                            # Try different encodings
                            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
                            content = None
                            
                            for encoding in encodings:
                                try:
                                    with open(file_path, 'r', encoding=encoding) as f:
                                        content = f.read()
                                    break  # If successful, exit the encoding loop
                                except UnicodeDecodeError:
                                    continue
                            
                            if content:
                                try:
                                    # Parse HTML and extract text
                                    soup = BeautifulSoup(content, 'lxml')
                                    # Remove script and style elements
                                    for script in soup(["script", "style"]):
                                        script.decompose()
                                    # Get text content
                                    text = soup.get_text(separator='\n', strip=True)
                                    if text:
                                        documents.append(Document(
                                            page_content=text,
                                            metadata={
                                                "source": file.name,
                                                "type": "scorm_content",
                                                "file": filename,
                                                "path": os.path.relpath(file_path, temp_dir)
                                            }
                                        ))
                                except Exception as e:
                                    st.warning(f"Error parsing HTML in {filename}: {str(e)}")
                            else:
                                st.warning(f"Could not decode file {filename} with any supported encoding")
                        
                        # Process video files
                        elif file_ext in ['.mp4', '.webm', '.avi', '.mov']:
                            st.info(f"Processing video file: {filename}")
                            video_text = extract_video_text(file_path)
                            if video_text:
                                documents.append(Document(
                                    page_content=video_text,
                                    metadata={
                                        "source": file.name,
                                        "type": "scorm_video",
                                        "file": filename,
                                        "path": os.path.relpath(file_path, temp_dir)
                                    }
                                ))
                            
                            # Also check for transcript files
                            transcript_exts = ['.txt', '.transcript', '.vtt', '.srt']
                            for ext in transcript_exts:
                                transcript_path = os.path.splitext(file_path)[0] + ext
                                if os.path.exists(transcript_path):
                                    with open(transcript_path, 'r', encoding='utf-8') as f:
                                        transcript_text = f.read()
                                        if transcript_text.strip():
                                            documents.append(Document(
                                                page_content=transcript_text,
                                                metadata={
                                                    "source": file.name,
                                                    "type": "scorm_video_transcript",
                                                    "video_file": filename,
                                                    "transcript_file": os.path.basename(transcript_path),
                                                    "path": os.path.relpath(transcript_path, temp_dir)
                                                }
                                            ))
                        
                    except Exception as e:
                        st.warning(f"Error processing file {filename}: {str(e)}")
    
    except Exception as e:
        st.error(f"Error processing SCORM package: {str(e)}")
    
    if documents:
        st.success(f"Successfully extracted content from {len(documents)} files in SCORM package")
    else:
        st.warning("No content could be extracted from the SCORM package")
    
    return documents 